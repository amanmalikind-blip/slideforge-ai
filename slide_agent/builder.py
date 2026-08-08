"""Render a DeckContent into a .pptx.

Two modes:
- Theme mode  : draws every slide from scratch on a 16:9 canvas using a built-in Theme.
- Template mode: opens the user's .pptx/.potx, keeps its masters/layouts/branding, fills
  native placeholders where possible and draws the rest using colours extracted from the file.
"""
from __future__ import annotations

import io
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .models import DeckContent, SlideContent
from .themes import DEFAULT_THEME_KEY, THEMES, Theme, theme_from_extracted
from . import template_analyzer as ta

# 16:9 canvas geometry (inches)
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.55
CONTENT_W = SLIDE_W - 2 * MARGIN


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _fmt(run, size: int, color: str, font: str, bold: bool = False, italic: bool = False):
    f = run.font
    f.size = Pt(size)
    f.name = font
    f.bold = bold
    f.italic = italic
    f.color.rgb = _rgb(color)


def _textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _para(tf, first: bool):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def _text(slide, x, y, w, h, text, size, color, font, bold=False, italic=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb, tf = _textbox(slide, x, y, w, h)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    _fmt(run, size, color, font, bold=bold, italic=italic)
    return tb


def _rect(slide, x, y, w, h, fill: str, line: Optional[str] = None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = _rgb(fill)
    if line:
        sp.line.color.rgb = _rgb(line)
        sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _bullets_into_tf(tf, bullets: List[str], theme: Theme, size: int):
    """Render bullets (with '>>' sub-bullet convention) into an existing text frame."""
    tf.word_wrap = True
    first = True
    for raw in bullets[:8]:
        text = str(raw).strip()
        if not text:
            continue
        sub = text.startswith(">>")
        if sub:
            text = text[2:].strip()
        p = _para(tf, first)
        first = False
        p.space_after = Pt(8 if not sub else 5)
        p.line_spacing = 1.12
        if sub:
            p.level = 1
            marker, msz, mcol = "–  ", size - 3, theme.muted
            body_col, bsz = theme.muted, size - 3
        else:
            marker, msz, mcol = "▪  ", size - 2, theme.accent
            body_col, bsz = theme.text, size
        r1 = p.add_run(); r1.text = marker; _fmt(r1, msz, mcol, theme.body_font, bold=not sub)
        r2 = p.add_run(); r2.text = text; _fmt(r2, bsz, body_col, theme.body_font)


def _bullet_size(bullets: List[str]) -> int:
    n = len(bullets)
    longest = max((len(b) for b in bullets), default=0)
    if n > 6 or longest > 120:
        return 15
    if n > 4 or longest > 90:
        return 16
    return 18


def _chrome(slide, theme: Theme, page_no: Optional[int], footer: str):
    """Footer text + page number on content slides."""
    if footer:
        _text(slide, MARGIN, SLIDE_H - 0.42, 6.0, 0.3, footer, 9, theme.muted, theme.body_font)
    if page_no:
        _text(slide, SLIDE_W - 1.1, SLIDE_H - 0.42, 0.6, 0.3, str(page_no), 10,
              theme.muted, theme.body_font, align=PP_ALIGN.RIGHT)


def _title_block(slide, theme: Theme, title: str, subtitle: str = "") -> float:
    """Standard header for content slides. Returns the y where body content may start."""
    _text(slide, MARGIN, 0.42, CONTENT_W, 0.9, title or " ", 27, theme.text,
          theme.heading_font, bold=True)
    _rect(slide, MARGIN + 0.02, 1.22, 1.05, 0.075, theme.accent)
    y = 1.5
    if subtitle:
        _text(slide, MARGIN, 1.42, CONTENT_W, 0.4, subtitle, 14, theme.muted, theme.body_font, italic=True)
        y = 1.95
    return y


# --------------------------------------------------------------------------- slide painters
def _paint_title(slide, sc: SlideContent, theme: Theme, deck: DeckContent, footer: str):
    bg = theme.bg if theme.is_dark else theme.primary
    fg = theme.text if theme.is_dark else "FFFFFF"
    sub = theme.muted if theme.is_dark else "DDE3F0"
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg)
    # decorative circles bleeding off the right edge
    _rect(slide, SLIDE_W - 3.4, -1.6, 5.2, 5.2, theme.secondary, shape=MSO_SHAPE.OVAL)
    _rect(slide, SLIDE_W - 1.7, 3.9, 3.4, 3.4, theme.accent, shape=MSO_SHAPE.OVAL)
    _rect(slide, MARGIN, 2.35, 1.3, 0.09, theme.accent)
    _text(slide, MARGIN, 2.6, 10.2, 1.9, sc.title or deck.title, 40, fg,
          theme.heading_font, bold=True, line_spacing=1.02)
    if sc.subtitle or deck.subtitle:
        _text(slide, MARGIN, 4.45, 9.6, 1.0, sc.subtitle or deck.subtitle, 19, sub, theme.body_font)
    if footer:
        _text(slide, MARGIN, SLIDE_H - 0.65, 8.0, 0.4, footer, 11, sub, theme.body_font)


def _paint_section(slide, sc: SlideContent, theme: Theme, index_label: str, footer: str):
    bg = theme.bg if theme.is_dark else theme.primary
    fg = theme.text if theme.is_dark else "FFFFFF"
    sub = theme.muted if theme.is_dark else "DDE3F0"
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg)
    _rect(slide, 0, 0, 0.35, SLIDE_H, theme.accent)
    _text(slide, MARGIN + 0.2, 1.7, 3.0, 1.6, index_label, 64, theme.accent,
          theme.heading_font, bold=True)
    _text(slide, MARGIN + 0.2, 3.15, 11.5, 1.5, sc.title, 33, fg, theme.heading_font,
          bold=True, line_spacing=1.05)
    if sc.subtitle:
        _text(slide, MARGIN + 0.2, 4.6, 10.5, 0.9, sc.subtitle, 16, sub, theme.body_font)
    if footer:
        _text(slide, MARGIN + 0.2, SLIDE_H - 0.6, 8.0, 0.4, footer, 10, sub, theme.body_font)


def _paint_bullets(slide, sc: SlideContent, theme: Theme):
    y = _title_block(slide, theme, sc.title, sc.subtitle)
    size = _bullet_size(sc.bullets)
    _, tf = _textbox(slide, MARGIN + 0.05, y + 0.15, CONTENT_W - 0.4, SLIDE_H - y - 1.0)
    _bullets_into_tf(tf, sc.bullets, theme, size)


def _paint_columns(slide, sc: SlideContent, theme: Theme, versus: bool):
    y = _title_block(slide, theme, sc.title, sc.subtitle)
    gap = 0.35
    col_w = (CONTENT_W - gap) / 2
    header_h = 0.55
    card_h = SLIDE_H - y - 1.05
    heads = [sc.left_title or ("Option A" if versus else ""), sc.right_title or ("Option B" if versus else "")]
    cols = [sc.left_bullets, sc.right_bullets]
    head_fills = [theme.primary, theme.accent] if versus else [theme.surface, theme.surface]
    head_colors = ["FFFFFF", "FFFFFF"] if versus else [theme.text, theme.text]
    if versus and theme.is_dark:
        head_colors = [theme.bg, theme.bg]
    for i in range(2):
        x = MARGIN + i * (col_w + gap)
        card = _rect(slide, x, y + 0.1, col_w, card_h, theme.surface, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.045
        if heads[i]:
            band = _rect(slide, x, y + 0.1, col_w, header_h, head_fills[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            band.adjustments[0] = 0.16
            btf = band.text_frame
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            btf.margin_left = Inches(0.18)
            p = btf.paragraphs[0]
            r = p.add_run(); r.text = heads[i]
            _fmt(r, 15, head_colors[i], theme.heading_font, bold=True)
        _, tf = _textbox(slide, x + 0.22, y + 0.28 + header_h, col_w - 0.44, card_h - header_h - 0.35)
        _bullets_into_tf(tf, cols[i][:5], theme, 14)


def _paint_quote(slide, sc: SlideContent, theme: Theme):
    _rect(slide, 0, 0, 0.35, SLIDE_H, theme.accent)
    _text(slide, 1.0, 0.7, 2.2, 1.8, "“", 110, theme.accent, "Georgia", bold=True)
    _text(slide, 1.7, 2.15, 10.0, 2.6, sc.quote or sc.title, 26, theme.text, theme.heading_font,
          italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    if sc.attribution:
        _text(slide, 1.7, 5.1, 10.0, 0.6, f"—  {sc.attribution}", 15, theme.muted,
              theme.body_font, align=PP_ALIGN.CENTER)


def _paint_kpi(slide, sc: SlideContent, theme: Theme):
    y = _title_block(slide, theme, sc.title, sc.subtitle)
    kpis = sc.kpis[:4] or []
    n = max(1, len(kpis))
    gap = 0.4
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = 2.3
    for i, k in enumerate(kpis):
        x = MARGIN + i * (card_w + gap)
        card = _rect(slide, x, y + 0.25, card_w, card_h, theme.surface, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.07
        _rect(slide, x, y + 0.25, card_w, 0.09, theme.accent)
        _text(slide, x + 0.1, y + 0.6, card_w - 0.2, 1.0, k.value, 34, theme.accent,
              theme.heading_font, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, x + 0.15, y + 1.65, card_w - 0.3, 0.8, k.label, 13, theme.muted,
              theme.body_font, align=PP_ALIGN.CENTER)
    if sc.bullets:
        _, tf = _textbox(slide, MARGIN + 0.05, y + 2.85, CONTENT_W - 0.4, 1.5)
        _bullets_into_tf(tf, sc.bullets[:3], theme, 14)


def _paint_process(slide, sc: SlideContent, theme: Theme):
    y = _title_block(slide, theme, sc.title, sc.subtitle)
    steps = [s for s in sc.steps if str(s).strip()][:6] or ["Step 1", "Step 2", "Step 3"]
    n = len(steps)
    overlap = 0.16
    w = (CONTENT_W + overlap * (n - 1)) / n
    h = 1.55
    top = y + 0.55
    palette = [theme.primary, theme.secondary, theme.accent]
    size = 13 if n <= 4 else 11
    for i, step in enumerate(steps):
        x = MARGIN + i * (w - overlap)
        chev = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(top), Inches(w), Inches(h))
        chev.fill.solid()
        chev.fill.fore_color.rgb = _rgb(palette[i % 3])
        chev.line.fill.background()
        chev.shadow.inherit = False
        tf = chev.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.28 if i else 0.16)
        tf.margin_right = Inches(0.22)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(step)
        _fmt(r, size, theme.bg if theme.is_dark else "FFFFFF", theme.body_font, bold=True)
        _text(slide, x + w / 2 - 0.3, top + h + 0.12, 0.6, 0.4, f"{i + 1:02d}", 13,
              theme.muted, theme.heading_font, bold=True, align=PP_ALIGN.CENTER)
    if sc.bullets:
        _, tf = _textbox(slide, MARGIN + 0.05, top + h + 0.7, CONTENT_W - 0.4, 1.8)
        _bullets_into_tf(tf, sc.bullets[:3], theme, 14)


def _paint_table(slide, sc: SlideContent, theme: Theme):
    y = _title_block(slide, theme, sc.title, sc.subtitle)
    spec = sc.table
    if not spec or not spec.headers:
        _paint_bullets(slide, sc, theme)
        return
    headers = [str(h) for h in spec.headers[:6]]
    rows = [[str(c) for c in r[:len(headers)]] for r in spec.rows[:8]]
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = min(0.62 * n_rows, SLIDE_H - y - 1.0)
    frame = slide.shapes.add_table(n_rows, n_cols, Inches(MARGIN), Inches(y + 0.15),
                                   Inches(CONTENT_W), Inches(height))
    table = frame.table
    table.first_row = True
    table.horz_banding = True
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(theme.primary)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = head
        _fmt(r, 13, "FFFFFF" if not theme.is_dark else theme.bg, theme.heading_font, bold=True)
    for ri, row in enumerate(rows, start=1):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(theme.surface if ri % 2 else theme.bg)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = row[ci] if ci < len(row) else ""
            _fmt(r, 12, theme.text, theme.body_font, bold=(ci == 0))


def _paint_closing(slide, sc: SlideContent, theme: Theme, footer: str):
    bg = theme.bg if theme.is_dark else theme.primary
    fg = theme.text if theme.is_dark else "FFFFFF"
    sub = theme.muted if theme.is_dark else "DDE3F0"
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg)
    _rect(slide, -1.9, 4.6, 4.4, 4.4, theme.secondary, shape=MSO_SHAPE.OVAL)
    _rect(slide, MARGIN, 2.2, 1.3, 0.09, theme.accent)
    _text(slide, MARGIN, 2.5, 12.0, 1.4, sc.title or "Thank you", 38, fg,
          theme.heading_font, bold=True)
    if sc.subtitle:
        _text(slide, MARGIN, 3.95, 10.5, 0.8, sc.subtitle, 18, sub, theme.body_font)
    if sc.bullets:
        _, tf = _textbox(slide, MARGIN, 4.8, 9.5, 1.8)
        for i, b in enumerate(sc.bullets[:3]):
            p = _para(tf, i == 0)
            p.space_after = Pt(6)
            r1 = p.add_run(); r1.text = "→  "; _fmt(r1, 15, theme.accent, theme.body_font, bold=True)
            r2 = p.add_run(); r2.text = str(b).lstrip(">").strip(); _fmt(r2, 15, fg, theme.body_font)
    if footer:
        _text(slide, MARGIN, SLIDE_H - 0.65, 9.0, 0.4, footer, 11, sub, theme.body_font)


_DRAWN_TYPES = {"quote", "kpi", "process", "table", "two_column", "comparison"}


def _paint(slide, sc: SlideContent, theme: Theme, deck: DeckContent, page_no: int,
           section_no: int, footer: str):
    kind = sc.type
    if kind == "title":
        _paint_title(slide, sc, theme, deck, footer)
    elif kind == "section":
        _paint_section(slide, sc, theme, f"{section_no:02d}", footer)
    elif kind == "two_column":
        _paint_columns(slide, sc, theme, versus=False); _chrome(slide, theme, page_no, footer)
    elif kind == "comparison":
        _paint_columns(slide, sc, theme, versus=True); _chrome(slide, theme, page_no, footer)
    elif kind == "quote":
        _paint_quote(slide, sc, theme); _chrome(slide, theme, page_no, footer)
    elif kind == "kpi":
        _paint_kpi(slide, sc, theme); _chrome(slide, theme, page_no, footer)
    elif kind == "process":
        _paint_process(slide, sc, theme); _chrome(slide, theme, page_no, footer)
    elif kind == "table":
        _paint_table(slide, sc, theme); _chrome(slide, theme, page_no, footer)
    elif kind == "closing":
        _paint_closing(slide, sc, theme, footer)
    else:
        _paint_bullets(slide, sc, theme); _chrome(slide, theme, page_no, footer)


def _add_notes(slide, notes: str):
    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass


# --------------------------------------------------------------------------- public API
def build_deck(
    deck: DeckContent,
    theme: Optional[Theme] = None,
    template_bytes: Optional[bytes] = None,
    template_name: str = "",
    footer: str = "",
) -> bytes:
    """Render the deck and return .pptx bytes."""
    if template_bytes:
        try:
            return _build_into_template(deck, template_bytes, template_name, footer)
        except Exception:
            # A broken/odd template should never block the user — fall back to theme mode.
            pass
    return _build_with_theme(deck, theme or THEMES[DEFAULT_THEME_KEY], footer)


def _build_with_theme(deck: DeckContent, theme: Theme, footer: str) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(int(Inches(SLIDE_W)))
    prs.slide_height = Emu(int(Inches(SLIDE_H)))
    blank = prs.slide_layouts[6]
    section_no = 0
    for i, sc in enumerate(deck.slides, start=1):
        slide = prs.slides.add_slide(blank)
        if sc.type != "title":  # paint background for content slides
            _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.bg).line.fill.background()
        if sc.type == "section":
            section_no += 1
        _paint(slide, sc, theme, deck, i, section_no, footer)
        _add_notes(slide, sc.notes)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _find_placeholder(slide, *ph_types):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type in ph_types and ph.has_text_frame:
                return ph
        except Exception:
            continue
    return None


def _fill_ph(ph, text: str):
    if ph is not None and text:
        ph.text_frame.text = text
        return True
    return False


def _fill_body_bullets(ph, bullets: List[str]):
    tf = ph.text_frame
    tf.clear()
    wrote = False
    for i, raw in enumerate([b for b in bullets if str(b).strip()][:8]):
        text = str(raw).strip()
        level = 1 if text.startswith(">>") else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text[2:].strip() if level else text
        p.level = level
        wrote = True
    return wrote


def _build_into_template(deck: DeckContent, data: bytes, filename: str, footer: str) -> bytes:
    prs = ta.open_template(data, filename)
    ta.remove_all_slides(prs)
    colors, fonts = ta.extract_theme_parts(prs)
    theme = theme_from_extracted(colors, fonts)
    PH = PP_PLACEHOLDER

    # Respect the template's own canvas size when drawing.
    global SLIDE_W, SLIDE_H, CONTENT_W  # noqa: PLW0603 - simple module-level geometry
    old = (SLIDE_W, SLIDE_H, CONTENT_W)
    SLIDE_W = prs.slide_width / 914400
    SLIDE_H = prs.slide_height / 914400
    CONTENT_W = SLIDE_W - 2 * MARGIN
    try:
        section_no = 0
        for i, sc in enumerate(deck.slides, start=1):
            kind = sc.type
            if kind == "section":
                section_no += 1
            if kind in _DRAWN_TYPES:
                # Drawn archetypes: use the emptiest layout as a canvas, keep template background.
                slide = prs.slides.add_slide(ta.blank_layout(prs))
                _paint(slide, sc, theme, deck, i, section_no, footer)
            else:
                layout = ta.pick_layout(prs, kind)
                slide = prs.slides.add_slide(layout)
                title_ph = slide.shapes.title or _find_placeholder(slide, PH.TITLE, PH.CENTER_TITLE)
                if not _fill_ph(title_ph, sc.title or (deck.title if kind == "title" else "")):
                    _text(slide, MARGIN, 0.5, CONTENT_W, 1.0,
                          sc.title or deck.title, 30, theme.text, theme.heading_font, bold=True)
                sub_ph = _find_placeholder(slide, PH.SUBTITLE)
                subtitle = sc.subtitle or (deck.subtitle if kind == "title" else "")
                body_ph = _find_placeholder(slide, PH.BODY, PH.OBJECT, PH.VERTICAL_BODY)
                if kind in ("title", "section", "closing"):
                    if not _fill_ph(sub_ph, subtitle) and subtitle:
                        if not (body_ph and _fill_body_bullets(body_ph, [subtitle])):
                            _text(slide, MARGIN, 2.3, CONTENT_W - 1, 0.8, subtitle, 18,
                                  theme.muted, theme.body_font)
                    if kind == "closing" and sc.bullets and body_ph:
                        _fill_body_bullets(body_ph, sc.bullets)
                else:  # bullets
                    _fill_ph(sub_ph, subtitle)
                    if body_ph:
                        _fill_body_bullets(body_ph, sc.bullets)
                    else:
                        _, tf = _textbox(slide, MARGIN + 0.05, 1.8, CONTENT_W - 0.4, SLIDE_H - 2.7)
                        _bullets_into_tf(tf, sc.bullets, theme, _bullet_size(sc.bullets))
            _add_notes(slide, sc.notes)
        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()
    finally:
        SLIDE_W, SLIDE_H, CONTENT_W = old
