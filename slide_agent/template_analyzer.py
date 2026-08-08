"""Understand a user-supplied PowerPoint template (.pptx or .potx).

What it does:
- Converts .potx template packages to .pptx in-memory (python-pptx cannot open .potx directly).
- Extracts the OOXML colour scheme (dk1/lt1/dk2/lt2/accent1..6) and font scheme from theme1.xml.
- Scores the template's slide layouts so the builder can pick the best layout per slide archetype.
- Strips any existing slides so only masters/layouts/branding remain.
"""
from __future__ import annotations

import io
import zipfile
from typing import Dict, Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

_CT_PRESENTATION = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
_CT_TEMPLATE = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
_CT_SLIDESHOW = "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml"

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def prepare_template_bytes(data: bytes, filename: str = "") -> bytes:
    """Return bytes python-pptx can open. Rewrites .potx/.ppsx content types to .pptx."""
    name = (filename or "").lower()
    if name.endswith(".pptx"):
        return data
    try:
        src = zipfile.ZipFile(io.BytesIO(data))
        ct = src.read("[Content_Types].xml").decode("utf-8")
    except (zipfile.BadZipFile, KeyError) as e:
        raise ValueError("That file does not look like a PowerPoint package.") from e

    if _CT_TEMPLATE not in ct and _CT_SLIDESHOW not in ct:
        return data  # already a normal presentation package

    ct = ct.replace(_CT_TEMPLATE, _CT_PRESENTATION).replace(_CT_SLIDESHOW, _CT_PRESENTATION)
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            payload = ct.encode("utf-8") if item.filename == "[Content_Types].xml" else src.read(item.filename)
            dst.writestr(item, payload)
    return out.getvalue()


def open_template(data: bytes, filename: str = "") -> Presentation:
    return Presentation(io.BytesIO(prepare_template_bytes(data, filename)))


# --------------------------------------------------------------------------- theme extraction
def extract_theme_parts(prs: Presentation) -> tuple[Dict[str, str], Dict[str, str]]:
    """Pull (colors, fonts) out of the first slide master's theme part."""
    colors: Dict[str, str] = {}
    fonts: Dict[str, str] = {}
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(RT.THEME)
        root = etree.fromstring(theme_part.blob)
        scheme = root.find(f".//{{{_A_NS}}}clrScheme")
        if scheme is not None:
            for slot in scheme:
                slot_name = etree.QName(slot).localname  # dk1, lt1, accent1, ...
                srgb = slot.find(f"{{{_A_NS}}}srgbClr")
                sys = slot.find(f"{{{_A_NS}}}sysClr")
                if srgb is not None and srgb.get("val"):
                    colors[slot_name] = srgb.get("val").upper()
                elif sys is not None and sys.get("lastClr"):
                    colors[slot_name] = sys.get("lastClr").upper()
        font_scheme = root.find(f".//{{{_A_NS}}}fontScheme")
        if font_scheme is not None:
            major = font_scheme.find(f"{{{_A_NS}}}majorFont/{{{_A_NS}}}latin")
            minor = font_scheme.find(f"{{{_A_NS}}}minorFont/{{{_A_NS}}}latin")
            if major is not None and major.get("typeface"):
                fonts["major"] = major.get("typeface")
            if minor is not None and minor.get("typeface"):
                fonts["minor"] = minor.get("typeface")
    except Exception:
        pass  # fall back to defaults downstream
    return colors, fonts


# --------------------------------------------------------------------------- slide management
def remove_all_slides(prs: Presentation) -> int:
    """Delete every existing slide (people often upload a full deck as their 'template')."""
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 - no public API for this yet
    removed = 0
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(f"{{{_R_NS}}}id")
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)
        removed += 1
    return removed


# --------------------------------------------------------------------------- layout picking
def _layout_features(layout) -> tuple[set, str]:
    types = set()
    for ph in layout.placeholders:
        try:
            types.add(ph.placeholder_format.type)
        except Exception:
            continue
    return types, (layout.name or "").lower()


def pick_layout(prs: Presentation, kind: str):
    """Choose the most appropriate layout in the template for a slide archetype."""
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("Template has no slide layouts.")

    PH = PP_PLACEHOLDER
    best, best_score = layouts[0], -999
    for layout in layouts:
        types, name = _layout_features(layout)
        n_ph = len(types)
        score = 0
        if kind in ("title", "closing"):
            if PH.CENTER_TITLE in types: score += 6
            if PH.SUBTITLE in types: score += 3
            if "title slide" in name or name.strip() == "title": score += 5
            if kind == "closing" and any(w in name for w in ("closing", "thank", "end")): score += 6
            if PH.BODY in types or PH.OBJECT in types: score -= 2
        elif kind == "section":
            if "section" in name or "divider" in name: score += 8
            if PH.TITLE in types or PH.CENTER_TITLE in types: score += 2
            if PH.BODY in types or PH.OBJECT in types: score -= 1
        elif kind in ("bullets",):
            if (PH.TITLE in types) and (PH.BODY in types or PH.OBJECT in types): score += 6
            if "title and content" in name or "content" in name: score += 3
            score -= max(0, n_ph - 3)  # avoid busy multi-placeholder layouts
        elif kind in ("two_column", "comparison"):
            body_like = sum(1 for t in types if t in (PH.BODY, PH.OBJECT))
            if "two content" in name or "comparison" in name: score += 8
            if body_like >= 2: score += 4
            if PH.TITLE in types: score += 2
        else:  # quote, kpi, process, table -> prefer title-only or blank canvases
            if n_ph == 0 or "blank" in name: score += 5
            if PH.TITLE in types and (PH.BODY not in types and PH.OBJECT not in types): score += 4
            if "title only" in name: score += 4
            score -= max(0, n_ph - 2)
        if score > best_score:
            best, best_score = layout, score
    return best


def blank_layout(prs: Presentation):
    """The emptiest layout available (used as a drawing canvas)."""
    layouts = list(prs.slide_layouts)
    return min(layouts, key=lambda l: len(list(l.placeholders)))


def describe_template(prs: Presentation) -> Dict[str, object]:
    """Summary shown in the UI after upload."""
    colors, fonts = extract_theme_parts(prs)
    return {
        "masters": len(prs.slide_masters),
        "layouts": [l.name for l in prs.slide_layouts],
        "existing_slides": len(prs.slides._sldIdLst),  # noqa: SLF001
        "colors": colors,
        "fonts": fonts,
        "size_in": (round(prs.slide_width / 914400, 2), round(prs.slide_height / 914400, 2)),
    }
